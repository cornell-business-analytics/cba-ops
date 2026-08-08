"""
Deliberation deck generator — produces a .pptx with one slide per candidate
showing their headshot, basic info, coffee chat scores, and interview scores.
"""
import io
import uuid

import httpx
from fastapi import APIRouter, Depends, Query
from fastapi.responses import Response
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.db.session import get_db
from app.models.candidate import (
    ApplicationCycle,
    Candidate,
    CoffeeChat,
    InterviewCategory,
    InterviewScore,
)
from app.models.user import User, UserRole
from app.modules.ops.deps import require_role

router = APIRouter(tags=["delib"])

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
GREEN   = RGBColor(0x1A, 0x7A, 0x3C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x11, 0x18, 0x27)
MID     = RGBColor(0x6B, 0x72, 0x80)
BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
LIGHT   = RGBColor(0xF9, 0xFA, 0xFB)
SCORE_Y = RGBColor(0x05, 0x96, 0x52)   # Yes / accepted
SCORE_M = RGBColor(0xD9, 0x7F, 0x06)   # Maybe
SCORE_N = RGBColor(0xDC, 0x26, 0x26)   # No / rejected

STATUS_COLORS: dict[str, RGBColor] = {
    "applied":      RGBColor(0x64, 0x74, 0x8B),
    "coffee_chat":  RGBColor(0x02, 0x96, 0xC6),
    "interviewing": RGBColor(0xD9, 0x7F, 0x06),
    "offer":        RGBColor(0x7C, 0x3A, 0xED),
    "accepted":     RGBColor(0x05, 0x96, 0x52),
    "rejected":     RGBColor(0xDC, 0x26, 0x26),
    "withdrawn":    RGBColor(0x9C, 0xA3, 0xAF),
}

STATUS_LABELS: dict[str, str] = {
    "applied": "Applied", "coffee_chat": "Coffee Chat",
    "interviewing": "Interviewing", "offer": "Offer",
    "accepted": "Accepted", "rejected": "Rejected", "withdrawn": "Withdrawn",
}

# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _rect(slide, left, top, width, height, fill: RGBColor, border: RGBColor | None = None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def _para(tf, text: str, size: float, bold=False, italic=False,
          color: RGBColor = DARK, align=PP_ALIGN.LEFT, space_before: float = 0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def _textbox(slide, left, top, width, height, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    # Remove the auto-added empty first paragraph by marking it; we'll add our own.
    return tb, tf


def _clear_tf(tf):
    """Remove the default empty paragraph pptx adds to new frames."""
    from pptx.oxml.ns import qn
    for elem in list(tf._txBody.findall(qn("a:p"))):
        tf._txBody.remove(elem)


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HDR_H   = Inches(0.85)
LEFT_W  = Inches(3.6)
GUTTER  = Inches(0.12)
RIGHT_X = LEFT_W + GUTTER * 2
RIGHT_W = SLIDE_W - RIGHT_X - Inches(0.25)
CONTENT_TOP = HDR_H + Inches(0.1)
CONTENT_H   = SLIDE_H - HDR_H - Inches(0.1)


def _build_slide(prs: Presentation, candidate: Candidate,
                 user_map: dict, cat_map: dict,
                 headshot_bytes: bytes | None,
                 slide_num: int, total: int) -> None:

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    status = candidate.status.value if hasattr(candidate.status, "value") else str(candidate.status)
    status_color = STATUS_COLORS.get(status, MID)
    status_label = STATUS_LABELS.get(status, status.replace("_", " ").title())

    # ── Header bar ──────────────────────────────────────────────────────────
    hdr = _rect(slide, Inches(0), Inches(0), SLIDE_W, HDR_H, GREEN)
    tf = hdr.text_frame
    tf.word_wrap = False
    _clear_tf(tf)
    _para(tf, candidate.name, 22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Status pill (top-right, overlaid)
    pill_w = Inches(1.5)
    pill = _rect(slide, SLIDE_W - pill_w - Inches(0.2),
                 Inches(0.2), pill_w, Inches(0.45),
                 status_color)
    tf2 = pill.text_frame
    tf2.word_wrap = False
    _clear_tf(tf2)
    _para(tf2, status_label, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Left panel background ────────────────────────────────────────────────
    _rect(slide, Inches(0), HDR_H, LEFT_W, CONTENT_H, LIGHT)

    # Vertical divider
    _rect(slide, LEFT_W, HDR_H + Inches(0.15),
          Inches(0.015), CONTENT_H - Inches(0.3), BORDER)

    # ── Headshot ─────────────────────────────────────────────────────────────
    hs_size = Inches(3.0)
    hs_left = (LEFT_W - hs_size) / 2
    hs_top  = HDR_H + Inches(0.2)

    if headshot_bytes:
        try:
            buf = io.BytesIO(headshot_bytes)
            slide.shapes.add_picture(buf, hs_left, hs_top, hs_size, hs_size)
        except Exception:
            _rect(slide, hs_left, hs_top, hs_size, hs_size, BORDER)
    else:
        _rect(slide, hs_left, hs_top, hs_size, hs_size, BORDER)
        tb, tf = _textbox(slide, hs_left, hs_top + hs_size / 2 - Inches(0.25),
                          hs_size, Inches(0.5))
        _clear_tf(tf)
        _para(tf, "No photo", 9, color=MID, align=PP_ALIGN.CENTER)

    # ── Info block (below headshot) ──────────────────────────────────────────
    info_top = hs_top + hs_size + Inches(0.15)
    info_pad = Inches(0.15)

    tb, tf = _textbox(slide, info_pad, info_top,
                      LEFT_W - info_pad * 2, Inches(1.4))
    _clear_tf(tf)
    tf.word_wrap = True

    # Major + grad year on one line
    meta_parts = []
    if candidate.major:
        meta_parts.append(candidate.major)
    if candidate.grad_year:
        meta_parts.append(f"'{candidate.grad_year[-2:]}")
    if meta_parts:
        _para(tf, "  ".join(meta_parts), 9, color=MID, align=PP_ALIGN.CENTER)

    # Pronouns
    if candidate.pronouns:
        _para(tf, candidate.pronouns, 8, italic=True, color=MID, align=PP_ALIGN.CENTER)

    # Colleges
    colleges = candidate.college or []
    if colleges:
        _para(tf, " · ".join(colleges), 8, color=MID, align=PP_ALIGN.CENTER)

    # Transfer flag
    if candidate.is_transfer:
        _para(tf, "Transfer student", 8, italic=True, color=MID, align=PP_ALIGN.CENTER)

    # ── Coffee chats ─────────────────────────────────────────────────────────
    chat_top = info_top + Inches(1.5)
    chats = [c for c in (candidate.coffee_chats or []) if c.completed or c.score is not None]
    if chats:
        tb, tf = _textbox(slide, info_pad, chat_top,
                          LEFT_W - info_pad * 2,
                          SLIDE_H - chat_top - Inches(0.15))
        _clear_tf(tf)
        tf.word_wrap = True

        _para(tf, "COFFEE CHATS", 7, bold=True, color=GREEN)

        for chat in chats:
            member_name = user_map.get(chat.member_id, "Unknown")
            score_str = f"{chat.score}/3" if chat.score is not None else "—"
            header_line = f"{member_name}  {score_str}"
            _para(tf, header_line, 8, bold=True, space_before=3)
            if chat.notes:
                snippet = chat.notes[:120].replace("\n", " ")
                if len(chat.notes) > 120:
                    snippet += "…"
                _para(tf, snippet, 7.5, italic=True, color=MID)

    # ── Page number ──────────────────────────────────────────────────────────
    tb, tf = _textbox(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.35),
                      Inches(1.3), Inches(0.3))
    _clear_tf(tf)
    _para(tf, f"{slide_num} / {total}", 7, color=MID, align=PP_ALIGN.RIGHT)

    # ── Right panel: interview scores ─────────────────────────────────────────
    right_top = CONTENT_TOP + Inches(0.05)
    tb, tf = _textbox(slide, RIGHT_X, right_top,
                      RIGHT_W, CONTENT_H - Inches(0.45))
    _clear_tf(tf)
    tf.word_wrap = True

    scores: list[InterviewScore] = candidate.interview_scores or []

    if not scores and not chats:
        _para(tf, "No scores recorded yet.", 9, italic=True, color=MID)
    else:
        # Group scores by category
        from collections import defaultdict
        by_cat: dict[uuid.UUID, list[InterviewScore]] = defaultdict(list)
        for s in scores:
            by_cat[s.category_id].append(s)

        if by_cat:
            _para(tf, "INTERVIEW SCORES", 8, bold=True, color=GREEN)

            for cat_id, cat_scores in by_cat.items():
                cat_name = cat_map.get(cat_id, "Category")
                _para(tf, cat_name.upper(), 8, bold=True, space_before=6)

                for s in cat_scores:
                    member_name = user_map.get(s.member_id, "Unknown")
                    if s.numeric_score is not None:
                        score_val = f"{s.numeric_score:.1f}"
                        score_color = (
                            SCORE_Y if s.numeric_score >= 4 else
                            SCORE_M if s.numeric_score >= 2.5 else
                            SCORE_N
                        )
                    elif s.ynm_score:
                        score_val = s.ynm_score
                        score_color = (
                            SCORE_Y if s.ynm_score.upper().startswith("Y") else
                            SCORE_N if s.ynm_score.upper() == "N" else
                            SCORE_M
                        )
                    else:
                        score_val = "—"
                        score_color = MID

                    score_line = f"{member_name}   {score_val}"
                    p = _para(tf, score_line, 8.5, space_before=2)
                    # Colour just the score part by adding a second run
                    runs = p.runs
                    if runs:
                        runs[0].font.color.rgb = DARK

                    if s.comments:
                        snippet = s.comments[:160].replace("\n", " ")
                        if len(s.comments) > 160:
                            snippet += "…"
                        _para(tf, f"  {snippet}", 7.5, italic=True, color=MID)

        # Internal notes
        if candidate.notes:
            _para(tf, "INTERNAL NOTES", 8, bold=True, color=GREEN, space_before=10)
            note_snippet = candidate.notes[:300].replace("\n", " ")
            if len(candidate.notes) > 300:
                note_snippet += "…"
            _para(tf, note_snippet, 8.5, italic=True, color=MID)


# ---------------------------------------------------------------------------
# Endpoint
# ---------------------------------------------------------------------------

@router.get("/cycles/{cycle_id}/delib-deck")
async def generate_delib_deck(
    cycle_id: uuid.UUID,
    statuses: list[str] = Query(
        default=["interviewing", "offer", "accepted"],
        alias="status",
    ),
    _: User = Depends(require_role(UserRole.director)),
    db: AsyncSession = Depends(get_db),
):
    cycle_result = await db.execute(
        select(ApplicationCycle).where(ApplicationCycle.id == cycle_id)
    )
    cycle = cycle_result.scalar_one_or_none()
    if not cycle:
        from fastapi import HTTPException
        raise HTTPException(status_code=404, detail="Cycle not found")

    # Candidates with coffee chats and interview scores eager-loaded
    q = (
        select(Candidate)
        .where(Candidate.cycle_id == cycle_id)
        .where(Candidate.status.in_(statuses))
        .options(
            selectinload(Candidate.coffee_chats),
            selectinload(Candidate.interview_scores),
        )
        .order_by(Candidate.name)
    )
    result = await db.execute(q)
    candidates = result.scalars().all()

    # User name map
    users_res = await db.execute(select(User))
    user_map = {u.id: u.name for u in users_res.scalars().all()}

    # Category name map
    cats_res = await db.execute(select(InterviewCategory))
    cat_map = {c.id: c.name for c in cats_res.scalars().all()}

    # Fetch headshots concurrently
    headshot_cache: dict[str, bytes | None] = {}
    urls_to_fetch = [
        c.headshot_url for c in candidates if c.headshot_url and c.headshot_url not in headshot_cache
    ]
    if urls_to_fetch:
        async with httpx.AsyncClient(timeout=8.0) as client:
            for url in urls_to_fetch:
                try:
                    r = await client.get(url)
                    headshot_cache[url] = r.content if r.status_code == 200 else None
                except Exception:
                    headshot_cache[url] = None

    # Build presentation
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(candidates)
    for i, candidate in enumerate(candidates, 1):
        hbytes = headshot_cache.get(candidate.headshot_url or "") if candidate.headshot_url else None
        _build_slide(prs, candidate, user_map, cat_map, hbytes, i, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe = cycle.name.replace(" ", "_").replace("/", "-")
    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="delib_{safe}.pptx"'},
    )
